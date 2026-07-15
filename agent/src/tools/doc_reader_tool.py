"""Universal document reader: dispatches by file extension.

Supported formats:
  - PDF (.pdf) — pypdfium2 + OCR fallback for image pages
  - Word (.docx) — python-docx (paragraphs + table cells)
  - Excel (.xlsx/.xls) — pandas preview, all sheets
  - PowerPoint (.pptx) — python-pptx (slide text)
  - Images (.png/.jpg/.jpeg/.gif/.bmp/.webp/.tiff) — OCR
  - Plain text (.txt/.md/.log/.json/.yaml/.yml/.toml/.ini/.cfg/.csv/.tsv/
                .html/.xml/.rst/.sql/.sh and common source-code extensions)

All handlers return the same JSON envelope: status, file, format, char_count,
truncated, text. PDF/Excel add format-specific metadata (pages, sheets, ...).
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Callable

from src.agent.progress import emit_progress
from src.agent.tools import BaseTool
from src.security.scanner import with_security_warnings
from src.tools.path_utils import safe_document_path

_MAX_CHARS = 15000
_MIN_TEXT_PER_PAGE = 50
_ENCODING_FALLBACK = ("utf-8", "utf-8-sig", "gbk", "gb2312", "big5", "latin-1")

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif"}
_TEXT_EXTS = {
    # docs / structured
    ".txt", ".md", ".log", ".rst",
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".csv", ".tsv", ".html", ".htm", ".xml",
    # source code (best-effort, LLM can read raw)
    ".py", ".pyi", ".js", ".jsx", ".ts", ".tsx",
    ".go", ".rs", ".java", ".kt", ".swift",
    ".c", ".h", ".cpp", ".hpp", ".cc",
    ".rb", ".php", ".pl", ".lua",
    ".sh", ".bash", ".zsh", ".ps1", ".bat",
    ".sql", ".r", ".m",
    ".dockerfile", ".makefile", ".cmake",
}

from src.tools.ocr import get_ocr_engine, get_ocr_install_hint

_cached_ocr_engine = None
_cached_ocr_checked = False


def _get_ocr():
    """Return the configured OCR engine (cached), or None."""
    global _cached_ocr_engine, _cached_ocr_checked
    if not _cached_ocr_checked:
        _cached_ocr_engine = get_ocr_engine()
        _cached_ocr_checked = True
    return _cached_ocr_engine


def _ocr_available() -> bool:
    return _get_ocr() is not None


def _ocr_image_array(img) -> str:
    """Run OCR on a numpy image via the pluggable engine."""
    engine = _get_ocr()
    if engine is None:
        return ""
    return engine.recognize(img)


# ---------------- shared helpers ----------------

def _err(msg: str) -> str:
    return json.dumps({"status": "error", "error": msg}, ensure_ascii=False)


def _truncate(text: str) -> tuple[str, bool]:
    """Clip to _MAX_CHARS, return (text, was_truncated)."""
    if len(text) <= _MAX_CHARS:
        return text, False
    return text[:_MAX_CHARS] + f"\n\n... (truncated, total {len(text)} chars)", True


def _envelope(path: Path, fmt: str, text: str, **extra: Any) -> str:
    """Build the standard JSON response."""
    body, truncated = _truncate(text)
    payload: dict[str, Any] = {
        "status": "ok",
        "file": path.name,
        "format": fmt,
        "char_count": len(text),
        "truncated": truncated,
        "text": body,
    }
    payload.update(extra)
    payload = with_security_warnings(payload, fields=("text",))
    return json.dumps(payload, ensure_ascii=False)


# ---------------- PDF ----------------

def _parse_pages(pages_str: str, total: int) -> list[int]:
    """Parse '1-10' / '5' / '1,3,5-8' into zero-based indices."""
    out: list[int] = []
    for part in pages_str.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start, end = part.split("-", 1)
            s = max(int(start.strip()) - 1, 0)
            e = min(int(end.strip()), total)
            out.extend(range(s, e))
        elif part.isdigit():
            out.append(int(part) - 1)
    return sorted(set(out))


def _read_pdf(path: Path, pages: str, min_text_per_page: int = _MIN_TEXT_PER_PAGE) -> str:
    """Extract PDF text; OCR pages with too little text."""
    try:
        import pypdfium2 as pdfium  # type: ignore
    except ImportError:
        return _err("pypdfium2 not installed; cannot read PDF")

    doc = pdfium.PdfDocument(str(path))
    try:
        total_pages = len(doc)
        targets = _parse_pages(pages, total_pages) if pages.strip() else list(range(total_pages))
        total_targets = len(targets)
        chunks: list[str] = []
        ocr_pages = 0
        skipped_pages = 0
        for idx, i in enumerate(targets, start=1):
            if not 0 <= i < total_pages:
                continue
            page = doc[i]
            text = page.get_textpage().get_text_range().strip()
            if len(text) >= min_text_per_page:
                chunks.append(f"--- Page {i + 1} ---\n{text}")
                emit_progress(
                    "reading_pdf",
                    current=idx,
                    total=total_targets,
                    message=f"page {i + 1}/{total_pages}",
                )
                continue
            if not _ocr_available():
                skipped_pages += 1
                emit_progress(
                    "reading_pdf",
                    current=idx,
                    total=total_targets,
                    message=f"page {i + 1}/{total_pages} (skipped: no OCR)",
                )
                continue
            bitmap = page.render(scale=300 / 72)
            img = bitmap.to_numpy()
            ocr_text = _ocr_image_array(img)
            if ocr_text.strip():
                chunks.append(f"--- Page {i + 1} [OCR] ---\n{ocr_text}")
                ocr_pages += 1
            elif text:
                chunks.append(f"--- Page {i + 1} ---\n{text}")
            emit_progress(
                "reading_pdf",
                current=idx,
                total=total_targets,
                message=f"page {i + 1}/{total_pages} (OCR)" if ocr_text.strip() else f"page {i + 1}/{total_pages}",
            )
        full = "\n\n".join(chunks)
        if not full and skipped_pages > 0:
            engine = _get_ocr()
            hint = get_ocr_install_hint(engine)
            return _err(
                f"All {total_pages} page(s) are scanned/image pages with no "
                f"extractable text, and no OCR engine is available. {hint}"
            )
        engine = _get_ocr()

        # Compute OCR quality metrics
        pages_read = len(targets)
        text_density = len(full) / max(pages_read, 1)

        if ocr_pages == 0:
            quality_flag = "no_ocr_needed" if not skipped_pages else "no_ocr_engine"
        elif skipped_pages > 0:
            quality_flag = "degraded"
        else:
            quality_flag = "good"

        ocr_quality = {
            "ocr_pages": ocr_pages,
            "text_density": round(text_density, 1),
            "quality_flag": quality_flag,
        }

        return _envelope(
            path, "pdf", full,
            total_pages=total_pages,
            pages_read=len(targets),
            ocr_pages=ocr_pages,
            skipped_pages=skipped_pages,
            ocr_engine=engine.name if engine else None,
            ocr_quality=ocr_quality,
        )
    finally:
        doc.close()


# ---------------- DOCX ----------------

def _read_docx(path: Path) -> str:
    try:
        import docx  # type: ignore
    except ImportError:
        return _err("python-docx not installed; run: pip install python-docx")

    doc = docx.Document(str(path))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for t_idx, table in enumerate(doc.tables, start=1):
        parts.append(f"\n--- Table {t_idx} ---")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            parts.append(" | ".join(cells))
    return _envelope(path, "docx", "\n".join(parts), paragraphs=len(doc.paragraphs), tables=len(doc.tables))


# ---------------- Excel ----------------

def _read_excel(path: Path) -> str:
    try:
        import pandas as pd  # type: ignore
    except ImportError:
        return _err("pandas not installed; cannot read Excel")

    xls = pd.ExcelFile(path)
    parts: list[str] = []
    sheet_info: list[dict[str, Any]] = []
    total_sheets = len(xls.sheet_names)
    for idx, name in enumerate(xls.sheet_names, start=1):
        emit_progress(
            "reading_excel",
            current=idx,
            total=total_sheets,
            message=f"sheet {name}",
        )
        df = xls.parse(name, dtype=str)
        preview = df.head(100).to_string(index=False)
        parts.append(f"--- Sheet: {name} ({len(df)} rows × {len(df.columns)} cols) ---\n{preview}")
        sheet_info.append({"name": name, "rows": len(df), "cols": len(df.columns)})
    return _envelope(path, "excel", "\n\n".join(parts), sheets=sheet_info)


# ---------------- PPTX ----------------

def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return _err("python-pptx not installed; run: pip install python-pptx")

    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {idx} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
    return _envelope(path, "pptx", "\n".join(parts), slides=len(prs.slides))


# ---------------- Image OCR ----------------

def _read_image(path: Path) -> str:
    try:
        import numpy as np  # type: ignore
        from PIL import Image  # type: ignore
    except ImportError:
        return _err("Pillow + numpy required for image OCR")

    try:
        img = np.array(Image.open(path).convert("RGB"))
    except Exception as exc:
        return _err(f"Failed to open image: {exc}")

    if not _ocr_available():
        engine = _get_ocr()
        hint = get_ocr_install_hint(engine)
        return _err(
            f"This image requires OCR to extract text, but no OCR engine is "
            f"available. {hint}"
        )

    text = _ocr_image_array(img)
    engine = _get_ocr()
    if not text.strip():
        ocr_quality = {
            "ocr_pages": 0,
            "text_density": 0.0,
            "quality_flag": "degraded",
        }
        return _envelope(
            path, "image", "",
            ocr_engine=engine.name if engine else None,
            ocr_quality=ocr_quality,
            note="OCR returned no text (empty or unreadable image)",
        )
    # Image = 1 page; text_density is total chars (matches PDF chars/page).
    ocr_quality = {
        "ocr_pages": 1,
        "text_density": round(len(text) / 1, 1),
        "quality_flag": "good",
    }
    return _envelope(
        path, "image", text,
        ocr_engine=engine.name if engine else None,
        ocr_quality=ocr_quality,
    )


# ---------------- Plain text ----------------

def _read_text(path: Path) -> str:
    """Read a text-like file with encoding fallback."""
    data = path.read_bytes()
    last_err: Exception | None = None
    for enc in _ENCODING_FALLBACK:
        try:
            decoded = data.decode(enc)
            return _envelope(path, "text", decoded, encoding=enc, size=len(data))
        except UnicodeDecodeError as exc:
            last_err = exc
    return _err(f"Failed to decode file with any of {_ENCODING_FALLBACK}: {last_err}")


# ---------------- Dispatcher ----------------

_HANDLERS: dict[str, Callable[[Path], str]] = {
    ".docx": _read_docx,
    ".xlsx": _read_excel,
    ".xls": _read_excel,
    ".pptx": _read_pptx,
}


def read_document(file_path: str, pages: str = "", min_text_per_page: int = _MIN_TEXT_PER_PAGE) -> str:
    """Read any supported document; dispatch by extension.

    Args:
        file_path: Absolute path to the file.
        pages: Only used for PDF — e.g. "1-10", "5", "1,3,5-8"; empty = all.
        min_text_per_page: Minimum text chars to consider a page as text-extracted
            (below this threshold triggers OCR for PDFs). Default 50.

    Returns:
        JSON envelope: status, file, format, char_count, truncated, text,
        plus format-specific metadata (total_pages, sheets, etc.).
    """
    try:
        path = safe_document_path(file_path)
    except ValueError as exc:
        return _err(str(exc))
    if not path.exists():
        return _err(f"File not found: {file_path}")
    if not path.is_file():
        return _err(f"Not a file: {file_path}")

    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return _read_pdf(path, pages, min_text_per_page=min_text_per_page)
        if ext in _HANDLERS:
            return _HANDLERS[ext](path)
        if ext in _IMAGE_EXTS:
            return _read_image(path)
        if ext in _TEXT_EXTS or ext == "":
            return _read_text(path)
        # Unknown extension: best-effort text read
        return _read_text(path)
    except Exception as exc:
        return _err(f"{type(exc).__name__}: {exc}")


class DocReaderTool(BaseTool):
    """Universal document reader — PDF/Word/Excel/PowerPoint/images/text."""

    name = "read_document"
    description = (
        "Read a document of any common format: PDF, Word (.docx), Excel "
        "(.xlsx/.xls), PowerPoint (.pptx), images (OCR), or plain text "
        "(txt/md/json/yaml/csv/html/code files). Returns extracted text in "
        "a unified JSON envelope. For PDFs, accepts an optional `pages` range."
    )
    parameters = {
        "type": "object",
        "properties": {
            "file_path": {"type": "string", "description": "Absolute path to the file."},
            "pages": {
                "type": "string",
                "description": "PDF only: page range (e.g. '1-10', '5', '1,3,5-8'). Ignored for other formats.",
                "default": "",
            },
            "min_text_per_page": {
                "type": "integer",
                "description": "PDF only: minimum text chars per page before OCR is triggered. Default 50.",
                "default": 50,
            },
        },
        "required": ["file_path"],
    }
    repeatable = True

    def execute(self, **kwargs: Any) -> str:
        return read_document(
            kwargs["file_path"],
            kwargs.get("pages", ""),
            min_text_per_page=kwargs.get("min_text_per_page", _MIN_TEXT_PER_PAGE),
        )
